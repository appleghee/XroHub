# -*- coding: utf-8 -*-
"""Draft 5 slide HVT — 16:9, palette chuẩn, placeholder ảnh, kịch bản trong Notes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RED, RED2, GOLD, GOLDD = RGBColor.from_string('C8102E'), RGBColor.from_string('7A0C1E'), RGBColor.from_string('F2B705'), RGBColor.from_string('D99A00')
BRICK, NAVY, CREAM, INK = RGBColor.from_string('8C5A3C'), RGBColor.from_string('1B2A4A'), RGBColor.from_string('F7F2E7'), RGBColor.from_string('232A33')
DARK, PAPER, CREAMT = RGBColor.from_string('240A10'), RGBColor.from_string('4A2E1C'), RGBColor.from_string('E9E1CE')
GREY = RGBColor.from_string('6E6250')

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, dash=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
        if dash:
            try: sp.line.dash_style = dash
            except Exception: pass
    return sp

def text(s, x, y, w, h, runs, size=14, bold=False, color=INK, font='Be Vietnam Pro',
         align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4):
    """runs: str hoặc list các đoạn (mỗi đoạn: str hoặc (str, dict))"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [runs]
    for i, r in enumerate(runs):
        cfg = {}
        if isinstance(r, tuple): r, cfg = r
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get('align', align); p.line_spacing = cfg.get('spacing', spacing)
        p.space_after = Pt(cfg.get('space_after', space_after))
        run = p.add_run(); run.text = r
        f = run.font
        f.name = cfg.get('font', font); f.size = Pt(cfg.get('size', size))
        f.bold = cfg.get('bold', bold); f.italic = cfg.get('italic', italic)
        f.color.rgb = cfg.get('color', color)
    return tb

def photo_ph(s, x, y, w, h, t1, t2, dark=True):
    sp = rect(s, x, y, w, h, fill=None, line=(RGBColor.from_string('D8CBB0') if not dark else RGBColor.from_string('CFC4AB')), lw=2.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: sp.line.dash_style = 4  # dash
    except Exception: pass
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '🖼️  ' + t1
    r.font.name = 'Montserrat'; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = RED2 if not dark else GOLD
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    r2 = p2.add_run(); r2.text = t2
    r2.font.name = 'Be Vietnam Pro'; r2.font.size = Pt(10.5); r2.font.color.rgb = GREY
    return sp

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ================= SLIDE 1 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 9.2, -1.2, 5.6, 5.6, fill=RGBColor.from_string('2A3A5E'), shape=MSO_SHAPE.STAR_5_POINT)
rect(s, 0, 0, 0.16, 7.5, fill=RED2)
text(s, 0.85, 0.72, 8.5, 0.4, 'BÀI TẬP LỚN • LỊCH SỬ & HOẠT ĐỘNG TRẢI NGHIỆM', size=13, bold=True, color=GOLD, font='Montserrat')
text(s, 0.8, 1.15, 8.6, 1.3, 'HOÀNG VĂN THỤ', size=57, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat')
text(s, 0.85, 2.55, 5, 0.6, '1909 – 1944', size=27, bold=True, color=GOLD, font='Montserrat')
text(s, 0.85, 3.2, 7.6, 0.6, '"Ngọn đuốc bất diệt & mái trường mang tên Người"', size=18, italic=True, color=CREAM, font='Playfair Display')
rect(s, 0.9, 3.85, 1.1, 0.07, fill=RED)
text(s, 0.85, 5.35, 7.4, 1.0, [('Trường THPT Hoàng Văn Thụ — Uông Bí, Quảng Ninh', {'bold': True, 'color': RGBColor.from_string('FFFFFF')}), ('Nhóm: …………… · Lớp: …………… · Tháng 9/2026', {'color': CREAMT})], size=13, spacing=1.25)
bg = rect(s, 0.85, 6.45, 4.4, 0.5, fill=None, line=GOLD, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.85, 6.5, 4.4, 0.4, 'TIẾT 1 + TIẾT 2 • THỜI LƯỢNG 8 PHÚT', size=11.5, bold=True, color=GOLD, font='Montserrat', align=PP_ALIGN.CENTER)
photo_ph(s, 8.15, 1.1, 4.45, 4.9, '[ẢNH 1] CHÂN DUNG ĐỒNG CHÍ', 'Ảnh tư liệu chính thức — bạn sẽ gửi · viền vàng đồng 4px')
cap = rect(s, 8.15, 6.12, 4.45, 0.46, fill=RGBColor.from_string('5A0F1C'), line=GOLD, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 8.15, 6.16, 4.45, 0.38, '1939 — BÍ THƯ XỨ ỦY BẮC KỲ', size=12, bold=True, color=GOLD, font='Montserrat', align=PP_ALIGN.CENTER)
notes(s, 'SLIDE 1 — 45 giây — MC1. [CUE] Nhạc nền trang trọng, âm lượng nhỏ; đứng yên 2 giây trước khi nói. Lời thoại: "Kính thưa thầy cô và các bạn! Lịch sử đôi khi không nằm trong những trang sách dày — nó nằm trong một cái tên. Cái tên ấy hôm nay nằm trên cổng trường chúng ta, nằm trong tên gọi của mỗi chúng ta: học sinh Trường THPT Hoàng Văn Thụ. Nhưng đằng sau cái tên đó là một con người — một chiến sĩ cộng sản chỉ sống 35 năm, mà cả một đời không một bước lùi. Hôm nay, nhóm chúng em xin được kể câu chuyện của Người — và của ngôi trường mang tên Người." [CUE] chìa tay về màn hình, chuyển Slide 2.')

# ================= SLIDE 2 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 0, 0, 13.333, 1.07, fill=NAVY)
text(s, 0.7, 0.22, 9.6, 0.6, 'TIẾT 1 · ĐỒNG CHÍ HOÀNG VĂN THỤ (1909 – 1944)', size=23, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat')
rect(s, 9.9, 0.3, 2.9, 0.46, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 9.9, 0.34, 2.9, 0.38, 'NGƯỜI THỢ – NGƯỜI IN – NGƯỜI LÀM BÁO', size=10, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat', align=PP_ALIGN.CENTER)
rect(s, 1.05, 1.65, 0.055, 5.35, fill=RED)
tl = [
    ('1909', 'Sinh ra tại Lạng Sơn · người dân tộc Tày · học chữ Hán & Quốc ngữ'),
    ('TRƯỚC 1939', 'Thợ cơ khí xưởng Nam Hưng · thợ in chữ ngược trên đá (Ngũ sắc thạch ấn) · làm báo Châu Giang, Tranh đấu, Lao động'),
    ('1939', 'Bí thư Xứ ủy Bắc Kỳ'),
    ('1940 – 1941', 'Thường vụ Trung ương Đảng · phụ trách công tác Binh vận & xây dựng căn cứ địa cách mạng'),
    ('1943 – 1944', 'Bị bắt 25/8/1943 → hy sinh 24/5/1944 tại Tương Mai'),
]
ys = [1.62, 2.42, 3.62, 4.32, 5.52]
for (yr, ds), y in zip(tl, ys):
    rect(s, 0.95, y + 0.05, 0.26, 0.26, fill=GOLD, line=RED, lw=1.75, shape=MSO_SHAPE.OVAL)
    text(s, 1.45, y - 0.04, 5.3, 0.35, yr, size=15.5, bold=True, color=RED2, font='Montserrat')
    text(s, 1.45, y + 0.33, 5.15, 0.95, ds, size=11.5, color=NAVY)
c1 = rect(s, 6.95, 1.5, 5.85, 2.2, fill=RGBColor.from_string('FFFFFF'), line=RGBColor.from_string('E7DEC9'), lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 6.95, 1.5, 0.11, 2.2, fill=BRICK)
text(s, 7.25, 1.66, 5.4, 0.35, 'BA NGHỀ TAY TRƯỚC CÁCH MẠNG', size=13, bold=True, color=BRICK, font='Montserrat')
text(s, 7.25, 2.08, 5.45, 1.5, ['⚙️  Thợ cơ khí — xưởng Nam Hưng', '🖨️  Thợ in chữ ngược trên đá — nhà in Ngũ sắc thạch ấn', '📰  Người làm báo — Châu Giang · Tranh đấu · Lao động'], size=12.5, color=INK, spacing=1.15)
c2 = rect(s, 6.95, 3.95, 5.85, 1.85, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.25, 4.1, 5.3, 0.32, 'VŨ KHÍ CỦA NGƯỜI', size=11.5, bold=True, color=GOLD, font='Montserrat')
text(s, 7.25, 4.44, 5.35, 0.75, '"Đôi tay thợ, máy in và ngòi bút trở thành vũ khí cách mạng."', size=15, italic=True, color=CREAM, font='Playfair Display')
for i, chip in enumerate(['CƠ SỞ LONG CHÂU', 'MẠNG LƯỚI BÍ MẬT']):
    rect(s, 7.25 + i * 2.15, 5.22, 2.0, 0.4, fill=None, line=GOLD, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 7.25 + i * 2.15, 5.27, 2.0, 0.32, chip, size=9.5, bold=True, color=GOLD, font='Montserrat', align=PP_ALIGN.CENTER)
photo_ph(s, 6.95, 6.05, 5.85, 1.05, '[ẢNH 2 — TÙY CHỌN]', 'Ảnh xưởng Nam Hưng / nhà in Ngũ sắc thạch ấn / trang báo cũ — bạn sẽ gửi', dark=False)
notes(s, 'SLIDE 2 — 1 phút 45 giây — MC2. "Đồng chí Hoàng Văn Thụ sinh năm 1909, người dân tộc Tày ở Lạng Sơn, thuở nhỏ học cả chữ Hán lẫn Quốc ngữ. Ít ai ngờ rằng, trước khi trở thành nhà cách mạng chuyên nghiệp, đồng chí từng là một người thợ: thợ cơ khí ở xưởng Nam Hưng; thợ in — người khắc chữ ngược trên đá cho nhà in Ngũ sắc thạch ấn; và là người làm báo cho các tờ Châu Giang, Tranh đấu, Lao động." [chỉ mốc 1939] "Từ năm 1939, đồng chí đảm nhiệm cương vị Bí thư Xứ ủy Bắc Kỳ. Giai đoạn 1940–1941, đồng chí là Ủy viên Thường vụ Trung ương Đảng, trực tiếp phụ trách công tác Binh vận và xây dựng căn cứ địa cách mạng." "Một người thợ, một người thợ in, một người làm báo — trở thành nhà lãnh đạo của cả một miền Bắc Kỳ. Chính ở những nghề tay chân giản dị ấy, đồng chí học được cách tổ chức, cách kết nối con người — nền tảng của tài năng lãnh đạo mà chúng ta sẽ thấy ngay sau đây." [CUE] hạ giọng câu cuối, chuyển Slide 3 + hạ đèn, nhạc trầm.')

# ================= SLIDE 3 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=DARK)
text(s, 0, 0.28, 13.333, 0.35, 'ĐIỂM NHẤN · MỘT HIỆN VẬT — MỘT CÂU CHUYỆN', size=12.5, bold=True, color=GOLD, font='Montserrat', align=PP_ALIGN.CENTER)
text(s, 0, 0.62, 13.333, 0.6, 'XÀ LIM SỐ 1 — NƠI Ý CHÍ KHÔNG THỂ BIẾN MẤT', size=26, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat', align=PP_ALIGN.CENTER)
photo_ph(s, 0.6, 1.42, 4.75, 3.6, '[ẢNH 3] XÀ LIM SỐ 1 (CELLULE 1)', 'Nhà tù Hỏa Lò — ảnh bạn sẽ gửi · tối ảnh 30% để hòa vào nền')
text(s, 0.6, 5.1, 4.75, 0.4, 'Hiện vật: Xà lim số 1 — nơi giam đồng chí sau ngày 25/8/1943', size=10.5, italic=True, color=CREAMT, align=PP_ALIGN.CENTER)
card = rect(s, 5.75, 1.42, 7.0, 3.75, fill=CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 6.0, 1.42, 1.2, 0.9, '“', size=52, bold=True, color=GOLD, font='Montserrat')
text(s, 6.0, 2.2, 6.5, 1.75, '[Trích 2–4 câu bài thơ "Nhắn bạn" — chỉ dùng trích đoạn đã được xác thực trong tư liệu của trường — bạn sẽ gửi/nhận từ thầy cô]', size=16.5, italic=True, color=PAPER, font='Playfair Display', spacing=1.3)
rect(s, 6.0, 4.32, 6.5, 0.03, fill=RGBColor.from_string('E0D5BE'))
text(s, 6.0, 4.45, 6.5, 0.6, '— ĐỒNG CHÍ HOÀNG VĂN THỤ · VIẾT BẰNG GẠCH NON TẠI XÀ LIM SỐ 1 (1943 – 1944)', size=11.5, bold=True, color=BRICK, font='Montserrat')
steps = [('25/8/1943', 'BỊ THỰC DÂN PHÁP BẮT'), ('XÀ LIM SỐ 1', 'NHÀ TÙ HỎA LÒ'), ('GẠCH NON', 'VIẾT "NHẮN BẠN" LÊN VÁCH TÙ'), ('24/5/1944 · TƯƠNG MAI', 'HIÊN NGANG — TỪ CHỐI BỊT MẮT')]
x = 0.6
for i, (a, b_) in enumerate(steps):
    sp = rect(s, x, 5.75, 3.0, 1.25, fill=RGBColor.from_string('3A1119'), line=GOLDD, lw=1.2, shape=MSO_SHAPE.CHEVRON)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = a; r.font.name = 'Montserrat'; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = GOLD
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = b_; r2.font.name = 'Be Vietnam Pro'; r2.font.size = Pt(9.5); r2.font.color.rgb = CREAMT
    x += 3.13
notes(s, 'SLIDE 3 — 2 phút 15 giây — MC3 + cả nhóm. Nói chậm hơn 30%. "Thưa thầy cô và các bạn, nếu chỉ được chọn một hiện vật để kể về đồng chí, nhóm chúng em chọn một không gian nhỏ và tăm tối: Xà lim số 1 — Cellule 1 — của nhà tù Hỏa Lò. Đêm 25/8/1943, đồng chí Hoàng Văn Thụ bị thực dân Pháp bắt và bị giam ở đây." "Trong xà lim số 1 ấy, đồng chí không có giấy, không có bút. Nhưng Người có một viên gạch non. Bằng viên gạch non ướt ấy, đồng chí viết lên vách tù bài thơ — Nhắn bạn. Gạch non sẽ khô và vụn đi theo năm tháng, nhưng niềm tin vào ngày thắng lợi của Người thì không bao giờ vụn mất." [ĐỌC THƠ] cả nhóm đọc đồng thanh trích đoạn ĐÃ XÁC THỰC (chèn khi có tư liệu). "Sáng ngày 24/5/1944, tại bãi Tương Mai, trước họng súng của kẻ thù, khi bị yêu cầu bịt mắt, đồng chí đã hiên ngang từ chối. Người muốn nhìn thẳng vào những phát đạn ấy — như Người đã nhìn thẳng vào mọi gian khổ của cách mạng. Đồng chí hy sinh khi mới 35 tuổi." "Xiềng xích có thể khóa tay chân một con người — nhưng không bao giờ khóa được ý chí của một người cộng sản." [CUE] nhạc dừng 2 giây — mặc niệm 3 giây — câu chốt — chuyển Slide 4.')

# ================= SLIDE 4 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 0, 0, 13.333, 1.07, fill=NAVY)
text(s, 0.7, 0.22, 9.8, 0.6, 'TIẾT 2 · LỊCH SỬ NGÔI TRƯỜNG MANG TÊN NGƯỜI', size=23, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat')
rect(s, 10.9, 0.3, 1.9, 0.46, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 10.9, 0.34, 1.9, 0.38, '1979 → NAY', size=12, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat', align=PP_ALIGN.CENTER)
rect(s, 0.6, 3.18, 12.13, 0.06, fill=RED)
cards = [
    ('1979 – 1984', 'PHÂN HIỆU CẤP 3 UÔNG BÍ', ['Bắt đầu với 3 lớp', 'Phát triển lên 6 lớp'], RED),
    ('17/11/1984', 'THÀNH LẬP TRƯỜNG', ['Quyết định số 645/QĐ-UBND', 'Tên đầu: Trường Trung học kỹ thuật Hoàng Văn Thụ'], BRICK),
    ('1989 – 2007', 'THPT HOÀNG VĂN THỤ', ['Sáp nhập THCS Nguyễn Văn Cừ', 'Liên cấp lớn nhất tỉnh: 40 lớp, hơn 1.800 HS'], NAVY),
    ('2007 – NAY', 'PHÁT TRIỂN BỀN VỮNG', ['Tách riêng cấp 3', 'Chuẩn Quốc gia (4/2012)', 'Bằng khen Thủ tướng CP (2014)'], GOLDD),
]
x = 0.6
for yr, ti, bl, cc in cards:
    rect(s, x, 1.55, 2.92, 2.95, fill=RGBColor.from_string('FFFFFF'), line=RGBColor.from_string('E7DEC9'), lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 1.55, 2.92, 0.14, fill=cc)
    text(s, x + 0.2, 1.8, 2.6, 0.45, yr, size=17, bold=True, color=RED2, font='Montserrat')
    text(s, x + 0.2, 2.28, 2.6, 0.4, ti, size=10.5, bold=True, color=NAVY, font='Montserrat')
    runs = [('• ' + b, {'space_after': 3}) for b in bl]
    text(s, x + 0.2, 2.72, 2.6, 1.6, runs, size=10.5, color=INK, spacing=1.1)
    x += 3.07
photo_ph(s, 0.6, 4.75, 5.6, 1.5, '[ẢNH 5] TRƯỜNG XƯA', 'Ảnh tư liệu giai đoạn đầu — bạn sẽ gửi', dark=False)
photo_ph(s, 7.13, 4.75, 5.6, 1.5, '[ẢNH 6] TRƯỜNG HÔM NAY', 'Ảnh cổng/sân trường hiện tại — bạn sẽ gửi', dark=False)
pills = ['★ Chuẩn Quốc gia — 4/2012', '★ Bằng khen Thủ tướng Chính phủ — 2014', '★ 40 lớp · hơn 1.800 HS (thời kỳ liên cấp)']
px = 1.15
for ptext in pills:
    w = 3.6 if '40 lớp' in ptext else 3.4
    rect(s, px, 6.55, w, 0.5, fill=RGBColor.from_string('FFF6DC'), line=GOLDD, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, px, 6.62, w, 0.38, ptext, size=11, bold=True, color=RGBColor.from_string('7A5C00'), font='Montserrat', align=PP_ALIGN.CENTER)
    px += w + 0.35
notes(s, 'SLIDE 4 — 1 phút 45 giây — MC1. "Sự hy sinh ấy đã được tri ân bằng một cái tên trường học. Hành trình của ngôi trường chúng ta bắt đầu từ năm 1979 — khi được tách thành Phân hiệu Cấp 3 Uông Bí với chỉ 3 lớp, và chỉ sau 5 năm đã phát triển lên 6 lớp." [thẻ 2] "Đến ngày 17/11/1984, Quyết định số 645/QĐ-UBND chính thức thành lập Trường Trung học kỹ thuật Hoàng Văn Thụ — cái tên liệt sĩ lần đầu tiên được in vào lịch sử của một ngôi trường tại Uông Bí." [thẻ 3] "Giai đoạn 1989–2007, trường đổi tên thành THPT Hoàng Văn Thụ, sáp nhập THCS Nguyễn Văn Cừ và trở thành trường liên cấp lớn nhất tỉnh với 40 lớp và hơn 1.800 học sinh." [thẻ 4] "Từ 2007 đến nay, sau khi tách cấp 2, trường không ngừng vươn lên: tháng 4/2012 đạt chuẩn Quốc gia, và năm 2014 vinh dự nhận Bằng khen của Thủ tướng Chính phủ. Từ 3 lớp học năm 1979 đến ngôi trường chuẩn quốc gia hôm nay — đó chính là cách thế hệ sau giữ lửa cho tên tuổi Người." [CUE] chuyển Slide 5.')

# ================= SLIDE 5 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=CREAM)
rect(s, 0, 0, 13.333, 1.07, fill=NAVY)
text(s, 0.7, 0.22, 11.5, 0.6, 'TỔNG KẾT · BỐN PHẨM CHẤT — MỘT NGỌN LỬA', size=23, bold=True, color=RGBColor.from_string('FFFFFF'), font='Montserrat')
qs = [('01', 'KIÊN TRUNG, BẤT KHUẤT', 'Không gục ngã trước họng súng Tương Mai'),
      ('02', 'TẬN TỤY, TRÁCH NHIỆM', 'Trọn đời vì sự nghiệp cách mạng'),
      ('03', 'TÀI NĂNG TỔ CHỨC & LÃNH ĐẠO', 'Cơ sở Long Châu, mạng lưới bí mật'),
      ('04', 'YÊU NƯỚC & TIN TẤT THẮNG', '"Nhắn bạn" viết từ xà lim số 1')]
pos = [(0.6, 1.5), (4.25, 1.5), (0.6, 4.15), (4.25, 4.15)]
for (num, tt, sub), (qx, qy) in zip(qs, pos):
    rect(s, qx, qy, 3.5, 2.45, fill=RGBColor.from_string('FFFFFF'), line=RGBColor.from_string('E7DEC9'), lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, qx, qy, 0.1, 2.45, fill=GOLD)
    circ = rect(s, qx + 0.25, qy + 0.25, 0.5, 0.5, fill=RED, shape=MSO_SHAPE.OVAL)
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.name = 'Montserrat'; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = RGBColor.from_string('FFFFFF')
    text(s, qx + 0.25, qy + 0.92, 3.0, 0.75, tt, size=13.5, bold=True, color=NAVY, font='Montserrat')
    text(s, qx + 0.25, qy + 1.62, 3.0, 0.7, sub, size=10.5, color=GREY)
rect(s, 8.35, 1.5, 4.4, 5.1, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
st = rect(s, 12.0, 1.65, 0.55, 0.55, fill=RGBColor.from_string('3A4C74'), shape=MSO_SHAPE.STAR_5_POINT)
text(s, 8.7, 1.75, 3.8, 0.35, 'THÔNG ĐIỆP', size=12, bold=True, color=GOLD, font='Montserrat')
text(s, 8.7, 2.2, 3.75, 2.6, '"Từ xà lim số 1 đến phòng học đầy nắng — ngọn lửa năm 1944 vẫn cháy trong mái trường mang tên Người."', size=17, italic=True, color=CREAM, font='Playfair Display', spacing=1.35)
rect(s, 8.7, 5.7, 3.7, 0.62, fill=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 8.7, 5.82, 3.7, 0.42, 'XIN CẢM ƠN THẦY CÔ & CÁC BẠN! 🌻', size=13, bold=True, color=NAVY, font='Montserrat', align=PP_ALIGN.CENTER)
notes(s, 'SLIDE 5 — 1 phút 15 giây — MC2 + MC3. MC2: "Từ xà lim số 1 đến những phòng học đầy nắng hôm nay — khoảng cách ấy chính là câu trả lời cho câu hỏi: vì sao chúng ta học lịch sử? Bởi bốn phẩm chất của đồng chí — kiên trung bất khuất, tận tụy trách nhiệm, tài năng tổ chức lãnh đạo, và lòng yêu nước cùng niềm tin tất thắng — không nằm trong quá khứ. Đó là những phẩm chất một học sinh Hoàng Văn Thụ hôm nay sống được mỗi ngày: kiên định trước khó khăn, trách nhiệm với bài tập và lời hứa, biết tổ chức công việc và tập thể, và luôn tin vào điều tốt đẹp." MC3: "Chúng em kể lại câu chuyện của Người không chỉ để tri ân — mà để tự nhủ: hãy sống cho xứng đáng với cái tên in trên chiếc áo mình mặc. Xin chân thành cảm ơn thầy cô và các bạn đã lắng nghe!" [CUE] nhạc lớn dần rồi tắt; cả 4 thành viên bước ra hàng, cúi chào; giữ slide 5 trong Q&A.')

prs.core_properties.title = 'Hoàng Văn Thụ (1909-1944) & Trường THPT Hoàng Văn Thụ — 5 slide'
prs.core_properties.author = 'Nhóm học sinh THPT Hoàng Văn Thụ'
prs.save('Draft_5_slide_HVT.pptx')
print('Saved Draft_5_slide_HVT.pptx với', len(prs.slides.slides if hasattr(prs.slides,"slides") else prs.slides._sldIdLst), 'slide')
