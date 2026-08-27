#!/usr/bin/env python3
"""
Tạo PowerPoint chuyên nghiệp về Lịch sử Quảng Ninh
Theo nguyên tắc thiết kế từ Anthropic Skills
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Palette màu chuyên nghiệp (lấy cảm hứng từ biển và than - đặc trưng Quảng Ninh)
COLORS = {
    'primary': RGBColor(15, 52, 84),      # Deep ocean blue
    'secondary': RGBColor(46, 117, 182),   # Brighter blue
    'accent': RGBColor(230, 126, 34),      # Warm orange (than/mặt trời)
    'light': RGBColor(236, 240, 241),      # Light gray background
    'dark': RGBColor(44, 62, 80),          # Dark text
    'white': RGBColor(255, 255, 255),
    'success': RGBColor(39, 174, 96),      # Green for achievements
}

def create_professional_pptx():
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient effect (simulated with rectangle)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # Decorative wave/line element
    wave = slide.shapes.add_shape(
        MSO_SHAPE.WAVE,
        Inches(0), Inches(6.5), prs.slide_width, Inches(1.5)
    )
    wave.fill.solid()
    wave.fill.fore_color.rgb = COLORS['accent']
    wave.line.fill.background()
    wave.rotation = 180
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11.333), Inches(2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "LỊCH SỬ TỈNH QUẢNG NINH"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    title_frame.word_wrap = True
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.2), Inches(9.333), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Hành trình qua các thời kỳ lịch sử dân tộc"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['light']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    #
    # SLIDE 2: Thời tiền sử và sơ sử
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()
    
    # Title in header
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(10), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "1. THỜI TIỀN SỬ VÀ SƠ SỬ"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Content area with light background
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = COLORS['light']
    content_bg.line.fill.background()
    content_bg.line.color.rgb = COLORS['secondary']
    content_bg.line.width = Pt(2)
    
    # Content text
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(11.333), Inches(4)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Con người sinh sống từ hàng nghìn năm trước",
        "Dấu tích để lại tại các hang động và đảo trên vịnh Hạ Long",
        "Thuộc lãnh thổ Văn Lang và Âu Lạc",
        "Nơi cư trú của người Việt cổ"
    ]
    
    for i, point in enumerate(points):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = f"•  {point}"
        para.font.size = Pt(24)
        para.font.color.rgb = COLORS['dark']
        para.space_after = Pt(20)
        para.line_spacing = 1.3
    
    #
    # SLIDE 3: Thời Bắc thuộc (2 cột layout)
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['secondary']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(10), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "2. THỜI BẮC THUỘC (Thế kỷ II TCN – Thế kỷ X)"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Left column - Context
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(5.8), Inches(5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    left_box.line.fill.background()
    
    left_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.8), Inches(5.4), Inches(4.5)
    )
    left_frame = left_text.text_frame
    left_frame.word_wrap = True
    
    para = left_frame.paragraphs[0]
    para.text = "Bối cảnh:"
    para.font.size = Pt(26)
    para.font.bold = True
    para.font.color.rgb = COLORS['primary']
    
    para2 = left_frame.add_paragraph()
    para2.text = "\nQuảng Ninh nằm dưới sự cai trị của các triều đại phong kiến phương Bắc trong hơn 1000 năm."
    para2.font.size = Pt(22)
    para2.font.color.rgb = COLORS['dark']
    para2.line_spacing = 1.4
    
    # Right column - Resistance
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.5), Inches(5.8), Inches(5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['accent']
    right_box.line.fill.background()
    
    right_text = slide.shapes.add_textbox(
        Inches(7.2), Inches(1.8), Inches(5.4), Inches(4.5)
    )
    right_frame = right_text.text_frame
    right_frame.word_wrap = True
    
    para = right_frame.paragraphs[0]
    para.text = "Khởi nghĩa tiêu biểu:"
    para.font.size = Pt(26)
    para.font.bold = True
    para.font.color.rgb = COLORS['white']
    
    resistance_points = [
        "• Khởi nghĩa Hai Bà Trưng",
        "• Các phong trào đấu tranh khác",
        "• Tinh thần bất khuất của nhân dân"
    ]
    
    for i, point in enumerate(resistance_points):
        if i == 0:
            para = right_frame.add_paragraph()
        else:
            para = right_frame.add_paragraph()
        para.text = point
        para.font.size = Pt(22)
        para.font.color.rgb = COLORS['white']
        para.space_after = Pt(15)
    
    #
    # SLIDE 4: Thời phong kiến độc lập
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(11), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "3. THỜI PHONG KIẾN ĐỘC LẬP (Thế kỷ X – XIX)"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Main content box
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.2)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLORS['light']
    content_box.line.fill.background()
    
    text_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(11.333), Inches(4.7)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    milestones = [
        ("Vùng biên giới quan trọng", "Vừa phát triển kinh tế vừa giữ vai trò phòng thủ"),
        ("Năm 938 - Chiến thắng Bạch Đằng", "Ngô Quyền đánh bại quân Nam Hán, mở ra thời kỳ độc lập"),
        ("Năm 1288 - Trận Bạch Đằng", "Trần Hưng Đạo chỉ huy chiến thắng quân Nguyên–Mông"),
        ("Yên Tử - Trung tâm tâm linh", "Thiền phái Trúc Lâm do Trần Nhân Tông sáng lập")
    ]
    
    for i, (title, desc) in enumerate(milestones):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        
        para.text = title
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = COLORS['primary']
        para.space_after = Pt(8)
        
        desc_para = text_frame.add_paragraph()
        desc_para.text = f"   → {desc}"
        desc_para.font.size = Pt(20)
        desc_para.font.color.rgb = COLORS['dark']
        desc_para.space_after = Pt(25)
        desc_para.line_spacing = 1.3
    
    #
    # SLIDE 5: Thời Pháp thuộc
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header with accent color
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['accent']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(10), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "4. THỜI PHÁP THUỘC (1883–1945)"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Grid layout with 3 boxes
    box_positions = [
        (Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.5)),
        (Inches(4.7), Inches(1.5), Inches(3.8), Inches(2.5)),
        (Inches(8.9), Inches(1.5), Inches(3.8), Inches(2.5)),
        (Inches(2.6), Inches(4.3), Inches(8.1), Inches(2))
    ]
    
    box_contents = [
        ("Khai thác than", "Hòn Gai, Cẩm Phả, Uông Bí\nđược khai thác mạnh", COLORS['primary']),
        ("Công nhân mỏ", "Hình thành giai cấp công nhân\ntrở thành lực lượng cách mạng", COLORS['secondary']),
        ("Truyền thống", "\"Kỷ luật và Đồng tâm\"\n- bản sắc công nhân mỏ", COLORS['dark']),
        ("Tổng bãi công 1936", "Cuộc đấu tranh lịch sử góp phần hình thành truyền thống cách mạng vững chắc", COLORS['accent'])
    ]
    
    for i, (pos, content) in enumerate(zip(box_positions, box_contents)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pos[0], pos[1], pos[2], pos[3]
        )
        box.fill.solid()
        box.fill.fore_color.rgb = content[2]
        box.line.fill.background()
        
        text_box = slide.shapes.add_textbox(
            pos[0] + Inches(0.3), pos[1] + Inches(0.25), 
            pos[2] - Inches(0.6), pos[3] - Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        title, desc, _ = content
        para = text_frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(22)
        para.font.bold = True
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER
        
        desc_para = text_frame.add_paragraph()
        desc_para.text = f"\n{desc}"
        desc_para.font.size = Pt(18)
        desc_para.font.color.rgb = COLORS['white']
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.line_spacing = 1.3
    
    #
    # SLIDE 6: Thời kỳ kháng chiến
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background with primary color
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "5. THỜI KỲ KHÁNG CHIẾN (1945–1975)"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Central emblem-like circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5), Inches(2), Inches(4.333), Inches(4.333)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['white']
    circle.line.color.rgb = COLORS['accent']
    circle.line.width = Pt(4)
    
    circle_text = slide.shapes.add_textbox(
        Inches(4.8), Inches(2.5), Inches(3.733), Inches(3.5)
    )
    circle_frame = circle_text.text_frame
    circle_frame.word_wrap = True
    
    para = circle_frame.paragraphs[0]
    para.text = "Chống Pháp\n&\nChống Mỹ"
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = COLORS['primary']
    para.alignment = PP_ALIGN.CENTER
    
    # Bottom info
    info_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(11.333), Inches(0.8)
    )
    info_frame = info_box.text_frame
    info_para = info_frame.paragraphs[0]
    info_para.text = "Than Quảng Ninh đóng vai trò quan trọng trong phát triển kinh tế và phục vụ đất nước"
    info_para.font.size = Pt(20)
    info_para.font.color.rgb = COLORS['light']
    info_para.alignment = PP_ALIGN.CENTER
    
    #
    # SLIDE 7: Từ 1963 đến nay (Timeline style)
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['secondary']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(10), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "6. TỪ NĂM 1963 ĐẾN NAY"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Timeline line
    timeline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(2.5), Inches(11.333), Inches(0.15)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS['accent']
    timeline.line.fill.background()
    
    # Milestone 1963
    milestone1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.8), Inches(2.1), Inches(0.8), Inches(0.8)
    )
    milestone1.fill.solid()
    milestone1.fill.fore_color.rgb = COLORS['primary']
    milestone1.line.fill.background()
    
    mile1_text = slide.shapes.add_textbox(
        Inches(1.3), Inches(3), Inches(1.8), Inches(1.5)
    )
    mile1_frame = mile1_text.text_frame
    mile1_para = mile1_frame.paragraphs[0]
    mile1_para.text = "30/10/1963\nThành lập tỉnh"
    mile1_para.font.size = Pt(16)
    mile1_para.font.bold = True
    mile1_para.font.color.rgb = COLORS['dark']
    mile1_para.alignment = PP_ALIGN.CENTER
    
    # Development areas after 1986
    dev_areas = [
        ("Khai thác & chế biến than", "Công nghiệp mũi nhọn"),
        ("Du lịch Vịnh Hạ Long", "Di sản Thiên nhiên Thế giới"),
        ("Thương mại biên giới", "Cửa khẩu Móng Cái"),
        ("Công nghiệp & Dịch vụ", "Hạ tầng hiện đại")
    ]
    
    for i, (area, desc) in enumerate(dev_areas):
        x_pos = Inches(1.5 + i * 3.2)
        y_pos = Inches(4.5)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(2.8), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['light']
        box.line.color.rgb = COLORS['secondary']
        box.line.width = Pt(2)
        
        box_text = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_pos + Inches(0.25), 
            Inches(2.4), Inches(1.5)
        )
        box_frame = box_text.text_frame
        box_frame.word_wrap = True
        
        para = box_frame.paragraphs[0]
        para.text = area
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = COLORS['primary']
        para.alignment = PP_ALIGN.CENTER
        
        desc_para = box_frame.add_paragraph()
        desc_para.text = f"\n{desc}"
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = COLORS['dark']
        desc_para.alignment = PP_ALIGN.CENTER
    
    #
    # SLIDE 8: Ý nghĩa lịch sử
    ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background simulation
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()
    
    # Decorative elements
    deco1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, Inches(0.3), prs.slide_height
    )
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = COLORS['accent']
    deco1.line.fill.background()
    
    deco2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - Inches(0.3), 0, Inches(0.3), prs.slide_height
    )
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = COLORS['accent']
    deco2.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Ý NGHĨA LỊCH SỬ"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Key points with icons (simulated with colored boxes)
    points = [
        ("Vị trí chiến lược", "Quốc phòng - An ninh biên giới"),
        ("Cái nôi công nghiệp than", "Ngành khai thác than Việt Nam"),
        ("Di sản thiên nhiên", "Vịnh Hạ Long - UNESCO"),
        ("Văn hóa tâm linh", "Yên Tử - Thiền phái Trúc Lâm"),
        ("Động lực phát triển", "Kinh tế vùng và đất nước")
    ]
    
    for i, (title, desc) in enumerate(points):
        y_pos = Inches(1.8 + i * 1.1)
        
        # Number circle
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.5), y_pos, Inches(0.6), Inches(0.6)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLORS['accent']
        num_circle.line.fill.background()
        
        num_text = slide.shapes.add_textbox(
            Inches(1.65), y_pos + Inches(0.15), Inches(0.3), Inches(0.3)
        )
        num_frame = num_text.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = str(i + 1)
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = COLORS['white']
        num_para.alignment = PP_ALIGN.CENTER
        
        # Text
        text_box = slide.shapes.add_textbox(
            Inches(2.3), y_pos + Inches(0.1), Inches(9), Inches(0.8)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        para = text_frame.paragraphs[0]
        para.text = f"{title}: {desc}"
        para.font.size = Pt(22)
        para.font.color.rgb = COLORS['light']
        para.line_spacing = 1.3
    
    # Save presentation
    output_file = '/workspace/lich_su_quang_ninh_chuyen_nghiep_v2.pptx'
    prs.save(output_file)
    print(f"✅ Đã tạo file PowerPoint chuyên nghiệp: {output_file}")
    print(f"📊 Tổng số slide: 8")
    print(f"🎨 Thiết kế theo nguyên tắc Anthropic Skills")

if __name__ == "__main__":
    create_professional_pptx()
