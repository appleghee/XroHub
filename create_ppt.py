from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Tạo presentation
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height

# Màu sắc chuyên nghiệp
COLOR_PRIMARY = RGBColor(0, 51, 102)      # Xanh đậm
COLOR_SECONDARY = RGBColor(204, 85, 0)    # Cam đất
COLOR_ACCENT = RGBColor(0, 153, 153)      # Xanh ngọc
COLOR_LIGHT = RGBColor(245, 245, 245)     # Xám nhạt
COLOR_WHITE = RGBColor(255, 255, 255)     # Trắng
COLOR_TEXT = RGBColor(50, 50, 50)         # Xám đen

def add_title_slide(prs, title_text, subtitle_text):
    """Tạo slide tiêu đề chuyên nghiệp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background gradient effect (rectangle với màu chính)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, 0, slide_width, slide_height
    )
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = COLOR_PRIMARY
    bg_fill.transparency = 0.1
    bg_shape.line.fill.background()
    
    # Title text box
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), slide_width - Inches(1), Inches(2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.bold = True
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = COLOR_WHITE
    title_para.alignment = PP_ALIGN.CENTER
    title_frame.word_wrap = True
    
    # Subtitle text box
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), slide_width - Inches(2), Inches(1.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle_text
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLOR_LIGHT
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_frame.word_wrap = True
    
    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(3), slide_width - Inches(4), Inches(0.1)
    )
    line_fill = line.fill
    line_fill.solid()
    line_fill.fore_color.rgb = COLOR_SECONDARY
    line.line.fill.background()
    
    return slide

def add_content_slide(prs, title_text, content_items):
    """Tạo slide nội dung chuyên nghiệp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide_width, Inches(1.2)
    )
    header_fill = header.fill
    header_fill.solid()
    header_fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # Title in header
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.bold = True
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = COLOR_WHITE
    title_frame.word_wrap = True
    
    # Content area background
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.4), slide_width - Inches(0.6), slide_height - Inches(1.7)
    )
    content_bg_fill = content_bg.fill
    content_bg_fill.solid()
    content_bg_fill.fore_color.rgb = COLOR_LIGHT
    content_bg_fill.transparency = 0.5
    content_bg.line.fill.background()
    
    # Content text
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.6), slide_width - Inches(1.2), slide_height - Inches(2.2)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        para.text = "• " + item
        para.font.size = Pt(18)
        para.font.color.rgb = COLOR_TEXT
        para.space_after = Pt(12)
        para.line_spacing = 1.3
    
    # Accent bar at bottom
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, slide_height - Inches(0.3), slide_width, Inches(0.3)
    )
    accent_fill = accent_bar.fill
    accent_fill.solid()
    accent_fill.fore_color.rgb = COLOR_SECONDARY
    accent_bar.line.fill.background()
    
    return slide

def add_two_column_slide(prs, title_text, left_title, left_items, right_title, right_items):
    """Tạo slide 2 cột chuyên nghiệp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide_width, Inches(1.2)
    )
    header_fill = header.fill
    header_fill.solid()
    header_fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.bold = True
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = COLOR_WHITE
    title_frame.word_wrap = True
    
    # Left column background
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.4), (slide_width - Inches(1)) / 2 - Inches(0.2), slide_height - Inches(1.9)
    )
    left_bg_fill = left_bg.fill
    left_bg_fill.solid()
    left_bg_fill.fore_color.rgb = COLOR_LIGHT
    left_bg_fill.transparency = 0.5
    left_bg.line.fill.background()
    
    # Right column background
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        slide_width / 2 + Inches(0.2), Inches(1.4), (slide_width - Inches(1)) / 2 - Inches(0.2), slide_height - Inches(1.9)
    )
    right_bg_fill = right_bg.fill
    right_bg_fill.solid()
    right_bg_fill.fore_color.rgb = COLOR_LIGHT
    right_bg_fill.transparency = 0.5
    right_bg.line.fill.background()
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), (slide_width - Inches(1)) / 2 - Inches(0.4), Inches(0.5)
    )
    left_title_frame = left_title_box.text_frame
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.text = left_title
    left_title_para.font.bold = True
    left_title_para.font.size = Pt(20)
    left_title_para.font.color.rgb = COLOR_ACCENT
    
    # Left column content
    left_content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), (slide_width - Inches(1)) / 2 - Inches(0.4), slide_height - Inches(2.5)
    )
    left_frame = left_content_box.text_frame
    left_frame.word_wrap = True
    left_frame.clear()
    
    for i, item in enumerate(left_items):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = "• " + item
        para.font.size = Pt(16)
        para.font.color.rgb = COLOR_TEXT
        para.space_after = Pt(10)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(
        slide_width / 2 + Inches(0.4), Inches(1.5), (slide_width - Inches(1)) / 2 - Inches(0.4), Inches(0.5)
    )
    right_title_frame = right_title_box.text_frame
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.text = right_title
    right_title_para.font.bold = True
    right_title_para.font.size = Pt(20)
    right_title_para.font.color.rgb = COLOR_ACCENT
    
    # Right column content
    right_content_box = slide.shapes.add_textbox(
        slide_width / 2 + Inches(0.4), Inches(2), (slide_width - Inches(1)) / 2 - Inches(0.4), slide_height - Inches(2.5)
    )
    right_frame = right_content_box.text_frame
    right_frame.word_wrap = True
    right_frame.clear()
    
    for i, item in enumerate(right_items):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = "• " + item
        para.font.size = Pt(16)
        para.font.color.rgb = COLOR_TEXT
        para.space_after = Pt(10)
    
    # Accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, slide_height - Inches(0.3), slide_width, Inches(0.3)
    )
    accent_fill = accent_bar.fill
    accent_fill.solid()
    accent_fill.fore_color.rgb = COLOR_SECONDARY
    accent_bar.line.fill.background()
    
    return slide

# Tạo slides
# Slide 1: Title
add_title_slide(prs, "LỊCH SỬ TỈNH QUẢNG NINH", "Hành trình qua các thời kỳ lịch sử dân tộc")

# Slide 2: Thời tiền sử và sơ sử
add_content_slide(prs, "1. THỜI TIỀN SỬ VÀ SƠ SỬ", [
    "Con người sinh sống từ hàng nghìn năm trước với dấu tích tại các hang động và đảo trên vịnh Hạ Long",
    "Thuộc lãnh thổ các nhà nước cổ Văn Lang và Âu Lạc",
    "Nơi cư trú của người Việt cổ"
])

# Slide 3: Thời Bắc thuộc
add_content_slide(prs, "2. THỜI BẮC THUỘC (Thế kỷ II TCN – Thế kỷ X)", [
    "Nằm dưới sự cai trị của các triều đại phong kiến phương Bắc",
    "Nhân dân tham gia nhiều cuộc khởi nghĩa chống Bắc thuộc",
    "Góp phần trong khởi nghĩa Hai Bà Trưng và các phong trào đấu tranh khác"
])

# Slide 4: Thời phong kiến độc lập
add_content_slide(prs, "3. THỜI PHONG KIẾN ĐỘC LẬP (Thế kỷ X – XIX)", [
    "Vùng biên giới quan trọng về kinh tế và quốc phòng",
    "Năm 938: Ngô Quyền đánh bại quân Nam Hán trên sông Bạch Đằng",
    "Năm 1288: Trần Hưng Đạo chiến thắng quân Nguyên–Mông tại Bạch Đằng",
    "Yên Tử trở thành trung tâm Thiền phái Trúc Lâm do Trần Nhân Tông sáng lập"
])

# Slide 5: Thời Pháp thuộc
add_content_slide(prs, "4. THỜI PHÁP THUỘC (1883–1945)", [
    "Thực dân Pháp khai thác mạnh các mỏ than ở Hòn Gai, Cẩm Phả, Uông Bí",
    "Hình thành giai cấp công nhân mỏ - lực lượng cách mạng quan trọng",
    "Năm 1936: Tổng bãi công của thợ mỏ",
    "Hình thành truyền thống \"Kỷ luật và Đồng tâm\""
])

# Slide 6: Thời kỳ kháng chiến
add_content_slide(prs, "5. THỜI KỲ KHÁNG CHIẾN (1945–1975)", [
    "Sau Cách mạng Tháng Tám 1945, tham gia kháng chiến chống Pháp và Mỹ",
    "Than Quảng Ninh đóng vai trò quan trọng trong phát triển kinh tế",
    "Phục vụ công cuộc xây dựng và bảo vệ đất nước"
])

# Slide 7: Từ 1963 đến nay (2 cột)
left_content = [
    "Ngày 30/10/1963: Thành lập tỉnh từ hợp nhất khu Hồng Quảng và tỉnh Hải Ninh",
    "Sau Đổi mới 1986, phát triển mạnh mẽ",
    "Khai thác và chế biến than",
    "Du lịch với Vịnh Hạ Long - Di sản Thiên nhiên Thế giới"
]

right_content = [
    "Thương mại biên giới với Trung Quốc qua Cửa khẩu Móng Cái",
    "Phát triển công nghiệp và dịch vụ",
    "Hạ tầng giao thông hiện đại",
    "Trở thành điểm đến du lịch hàng đầu Việt Nam"
]

add_two_column_slide(prs, "6. TỪ NĂM 1963 ĐẾN NAY", "Thành lập & Phát triển", left_content, "Các lĩnh vực trọng điểm", right_content)

# Slide 8: Ý nghĩa lịch sử
add_content_slide(prs, "Ý NGHĨA LỊCH SỬ", [
    "Vị trí chiến lược về quốc phòng-an ninh",
    "Cái nôi của ngành công nghiệp khai thác than Việt Nam",
    "Nổi tiếng với di sản thiên nhiên và văn hóa thế giới",
    "Đóng góp quan trọng vào phát triển kinh tế và lịch sử dân tộc"
])

# Lưu file
prs.save('lich_su_quang_ninh_chuyen_nghiep.pptx')
print("✅ Đã tạo PowerPoint chuyên nghiệp: lich_su_quang_ninh_chuyen_nghiep.pptx")
